import os
import base64
import io
from io import BytesIO
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
import csv
import openpyxl
from bs4 import BeautifulSoup
import markdown2
from PIL import Image

def read_document(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext in ['.docx', '.doc']:
            return read_word_document(file_path)
        elif ext == '.pdf':
            return read_pdf_document(file_path)
        elif ext in ['.pptx', '.ppt']:
            return read_powerpoint_document(file_path)
        elif ext == '.rtf':
            return read_rtf_document(file_path)
        elif ext == '.odt':
            return read_odt_document(file_path)
        elif ext in ['.txt', '.md', '.html', '.htm']:
            return read_text_document(file_path)
        else:
            raise ValueError(f"Unsupported document format: {ext}")
    except Exception as e:
        raise IOError(f"Error reading document {file_path}: {str(e)}")

def process_image(image_path_or_bytes, preserve_original=False):
    try:
        # Open the image from file path or bytes
        if isinstance(image_path_or_bytes, str):
            image = Image.open(image_path_or_bytes)
        else:
            image = Image.open(io.BytesIO(image_path_or_bytes))
        
        if not preserve_original:
            # Convert image based on its mode
            if image.mode in ['RGBA', 'LA']:
                # Images with alpha channel
                background = Image.new('RGBA', image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[3])  # 3 is the alpha channel
                image = background.convert('RGB')
            elif image.mode == 'P':
                # Palette images
                image = image.convert('RGBA')
                background = Image.new('RGBA', image.size, (255, 255, 255))
                background.paste(image, mask=image.split()[3])
                image = background.convert('RGB')
            elif image.mode != 'RGB':
                # All other modes (grayscale, CMYK, YCbCr, etc.)
                image = image.convert('RGB')

            # Resize the image if it's too large
            max_size = (1280, 720)  # 720p
            image.thumbnail(max_size, Image.LANCZOS)

        # Save the image to a byte stream
        buffered = io.BytesIO()
        image.save(buffered, format="JPEG", quality=85)
        
        # Encode the image to base64
        img_str = base64.b64encode(buffered.getvalue()).decode()
        
        return img_str
    except Exception as e:
        print(f"Error processing image: {str(e)}")
        return None

def read_pdf_document(file_path):
    doc = fitz.open(file_path)
    content = []
    images = []

    for page in doc:
        text = page.get_text("blocks")
        for block in text:
            if block[6] == 0:  # If it's a text block
                content.append(block[4].strip())
        
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            img_str = process_image(image_bytes)
            if img_str:
                images.append(img_str)

    return "\n\n".join(content), images

def read_word_document(file_path):
    doc = Document(file_path)
    content = []
    images = []

    for para in doc.paragraphs:
        if para.text.strip():
            content.append(para.text.strip())

    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_part = rel.target_part
            image_bytes = image_part.blob
            img_str = process_image(image_bytes)
            if img_str:
                images.append(img_str)

    return "\n\n".join(content), images

def read_powerpoint_document(file_path):
    prs = Presentation(file_path)
    content = []
    images = []
    slide_count = len(prs.slides)

    content.append(f"PowerPoint presentation with {slide_count} slides.")

    for i, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # Limit text to first 100 characters per shape
                slide_content.append(shape.text.strip()[:100])
            if shape.shape_type == 13:  # Picture
                image_bytes = shape.image.blob
                try:
                    # Use process_image function
                    processed_image = process_image(image_bytes)
                    if processed_image:
                        images.append(processed_image)
                        slide_content.append("[Image]")
                except Exception as e:
                    print(f"Error processing image in slide {i}: {str(e)}")

        # Summarize slide content
        if slide_content:
            summary = f"Slide {i}: " + " | ".join(slide_content)
            content.append(summary[:200])  # Limit each slide summary to 200 characters

        # Only process notes if they exist and are non-empty
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes = slide.notes_slide.notes_text_frame.text.strip()
            content.append(f"Notes for Slide {i}: {notes[:100]}")  # Limit notes to 100 characters

    # Limit total content to 2000 characters
    full_content = "\n".join(content)
    if len(full_content) > 2000:
        full_content = full_content[:1997] + "..."

    return full_content, images

def read_rtf_document(file_path):
    # RTF reading is complex and requires a dedicated library
    # For now, we'll just read it as plain text
    with open(file_path, 'r', errors='ignore') as file:
        content = file.read()
    return content, []

def read_odt_document(file_path):
    # ODT reading requires a dedicated library like odfpy
    # For now, we'll return an error message
    return "ODT file format is not supported yet.", []

def read_text_document(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()

    if ext == '.md':
        content = markdown2.markdown(content)
    elif ext in ['.html', '.htm']:
        soup = BeautifulSoup(content, 'html.parser')
        content = soup.get_text()

    return content, []

def read_spreadsheet(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.csv':
            with open(file_path, 'r', newline='', encoding='utf-8-sig') as csvfile:
                reader = csv.reader(csvfile)
                content = '\n'.join([', '.join(row) for row in reader])
        elif ext in ['.xlsx', '.xls']:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets_data = []
            for sheet in wb.worksheets:
                sheet_data = f"Sheet: {sheet.title}\n"
                sheet_data += '\n'.join([', '.join(str(cell.value) if cell.value is not None else '' for cell in row) for row in sheet.iter_rows()])
                sheets_data.append(sheet_data)
            content = '\n\n'.join(sheets_data)
        else:
            raise ValueError(f"Unsupported spreadsheet format: {ext}")

        return content, []
    except Exception as e:
        raise IOError(f"Error reading spreadsheet {file_path}: {str(e)}")